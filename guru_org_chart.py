from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Slide setup: 16:9 widescreen ───────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Colour palette ──────────────────────────────────────────────────────────
C_TOP   = RGBColor(0x1F, 0x38, 0x64)   # dark navy  – top node
C_L1    = RGBColor(0x2E, 0x75, 0xB6)   # mid blue   – direct reports
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY  = RGBColor(0x60, 0x60, 0x60)
C_LINE  = RGBColor(0x8E, 0xA9, 0xC8)

# ── Helper: styled rounded-rectangle box ───────────────────────────────────
def add_box(line1, line2, left, top, width=2.6, height=0.72,
            fill=C_L1, text_color=C_WHITE, size=9, bold=False):
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE_TYPE ROUNDED_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.adjustments[0] = 0.05

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = C_WHITE
    shape.line.width = Pt(0.75)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # First paragraph: department (or name if no dept)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = line1
    r1.font.size = Pt(size)
    r1.font.color.rgb = text_color
    r1.font.bold = True

    # Second paragraph: name (only if there is one)
    if line2:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = line2
        r2.font.size = Pt(size - 1)
        r2.font.color.rgb = RGBColor(0xD0, 0xE4, 0xFF)
        r2.font.bold = False

    return shape

# ── Helper: connector ───────────────────────────────────────────────────────
def connect(a, b):
    conn = slide.shapes.add_connector(
        1,
        a.left + a.width  // 2,
        a.top  + a.height,
        b.left + b.width  // 2,
        b.top,
    )
    conn.line.color.rgb = C_LINE
    conn.line.width = Pt(1.5)

# ── Layout constants ────────────────────────────────────────────────────────
W    = 13.33
Y0   = 1.20   # Guru row
Y1   = 3.10   # direct-report row

TW   = 2.4    # top box width
TH   = 0.78   # top box height

BW   = 2.6    # report box width
BH   = 0.80   # report box height
GAP  = 0.50   # gap between sibling boxes

# ── Slide subtitle ──────────────────────────────────────────────────────────
tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.28))
tf = tb.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run = tf.paragraphs[0].add_run()
run.text = "Organisation Chart — Guru"
run.font.size = Pt(9)
run.font.color.rgb = C_GREY
run.font.bold = True

# ── Top node: Guru ──────────────────────────────────────────────────────────
guru = add_box("Guru", "", W / 2 - TW / 2, Y0,
               width=TW, height=TH, fill=C_TOP, size=14, bold=True)

# ── Direct reports ──────────────────────────────────────────────────────────
# Structure: (department, name)  — if no dept, dept == name and name == ""
reports = [
    ("Service Ops & Revenue", "Sajeev"),
    ("Home Health",           "Vineeth Upadhyay"),
    ("Kuberan",               ""),          # no department, just name
]

n      = len(reports)
total  = n * BW + (n - 1) * GAP
start  = (W - total) / 2

boxes = []
for i, (dept, name) in enumerate(reports):
    bx = start + i * (BW + GAP)
    box = add_box(dept, name, bx, Y1, width=BW, height=BH)
    boxes.append(box)
    connect(guru, box)

prs.save("guru_org_chart.pptx")
print("Saved: guru_org_chart.pptx")
