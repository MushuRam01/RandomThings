from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Slide setup: standard widescreen 16:9 ──────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Colour palette ──────────────────────────────────────────────────────────
C_TOP   = RGBColor(0x1F, 0x38, 0x64)   # dark navy   – CEO
C_L1    = RGBColor(0x2E, 0x75, 0xB6)   # medium blue – direct reports
C_L2    = RGBColor(0x5B, 0x9B, 0xD5)   # sky blue    – second level
C_L3    = RGBColor(0xBD, 0xD7, 0xEE)   # pale blue   – third level
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
C_LINE  = RGBColor(0x8E, 0xA9, 0xC8)

# ── Helper: styled rounded-rectangle box ───────────────────────────────────
def add_box(text, left, top, width=1.4, height=0.56,
            fill=C_L1, text_color=C_WHITE, size=7.5, bold=False):
    shape = slide.shapes.add_shape(
        5,              # MSO_SHAPE_TYPE "ROUNDED_RECTANGLE"
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.adjustments[0] = 0.05   # small corner radius

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = C_WHITE
    shape.line.width = Pt(0.75)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = text_color
    run.font.bold = bold
    return shape

# ── Helper: connector line ──────────────────────────────────────────────────
def connect(a, b):
    cx_a = a.left + a.width  // 2
    cy_a = a.top  + a.height
    cx_b = b.left + b.width  // 2
    cy_b = b.top
    conn = slide.shapes.add_connector(1, cx_a, cy_a, cx_b, cy_b)
    conn.line.color.rgb = C_LINE
    conn.line.width = Pt(1.25)

# ── Layout constants (all in inches) ───────────────────────────────────────
W       = 13.33          # slide width
Y0      = 0.30           # CEO row
Y1      = 1.45           # direct-report row
Y2      = 2.80           # sub-report row
Y3      = 4.10           # third-level row

BW      = 1.38           # standard box width
BH      = 0.58           # standard box height
GAP     = 0.235          # horizontal gap between siblings
STEP    = BW + GAP

# ── Slide title (decorative) ────────────────────────────────────────────────
title = slide.shapes.add_textbox(Inches(0.3), Inches(0.0), Inches(12.7), Inches(0.28))
tf = title.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run = tf.paragraphs[0].add_run()
run.text = "Organisation Chart — Amit Nagpal"
run.font.size = Pt(9)
run.font.color.rgb = RGBColor(0x60, 0x60, 0x60)
run.font.bold = True

# ── CEO ─────────────────────────────────────────────────────────────────────
amit = add_box("Amit Nagpal", W/2 - 1.15, Y0,
               width=2.3, height=0.62, fill=C_TOP, size=11, bold=True)

# ── 8 direct reports (level 1) ──────────────────────────────────────────────
# Total occupied width: 8*BW + 7*GAP
l1_total = 8 * BW + 7 * GAP
l1_start = (W - l1_total) / 2

L1_LABELS = [
    "North + Nepal\nMandeep",
    "South\nGiri",
    "West\n—",
    "East\nAnirban",
    "Channel\nSales",
    "Govt\nRupesh",
    "International\nBusiness",
    "Pharma\nSales",
]

l1_boxes = []
for i, label in enumerate(L1_LABELS):
    bx = l1_start + i * STEP
    box = add_box(label, bx, Y1, width=BW, height=BH, fill=C_L1, size=7.5)
    l1_boxes.append(box)
    connect(amit, box)

# ── Level 2: Shashibhusan under Pharma ─────────────────────────────────────
pharma = l1_boxes[7]
pharma_cx = l1_start + 7 * STEP + BW / 2
shashi = add_box("Shashibhusan", pharma_cx - 0.8, Y2,
                 width=1.6, height=BH, fill=C_L2, size=7.5)
connect(pharma, shashi)

# ── Level 2: NCM under Channel ──────────────────────────────────────────────
channel = l1_boxes[4]
chan_cx = l1_start + 4 * STEP + BW / 2
NCM_W = 2.2
ncm = add_box("National Channel Sales Mgr\nAnzar Bhatt / Sharddha",
              chan_cx - NCM_W / 2, Y2,
              width=NCM_W, height=0.68, fill=C_L2, size=7.5)
connect(channel, ncm)

# ── Level 3: CN CS CW CE under NCM ─────────────────────────────────────────
L3W = 0.82
L3_GAP = 0.22
l3_total = 4 * L3W + 3 * L3_GAP
l3_start = chan_cx - l3_total / 2

for i, region in enumerate(["North", "South", "West", "East"]):
    bx = l3_start + i * (L3W + L3_GAP)
    box = add_box(region, bx, Y3, width=L3W, height=0.46,
                  fill=C_L3, text_color=C_DARK, size=7.5)
    connect(ncm, box)

prs.save("amit_nagpal_org_chart.pptx")
print("Saved: amit_nagpal_org_chart.pptx")